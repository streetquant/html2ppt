import argparse
import asyncio
import os
import re
import sys
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.action import PP_ACTION
import base64
from PIL import Image
import io

# --- CONFIGURATION ---
# Google Slides uses 13.33" x 7.5" at 96 DPI (1280x720) or 16:9 aspect ratio
SLIDE_WIDTH_PX = 1280  # Changed to 16:9 aspect ratio (Google Slides standard)
SLIDE_HEIGHT_PX = 720  # Changed to 16:9 aspect ratio (Google Slides standard)
PX_TO_INCH = 1/96
PPTX_WIDTH = Inches(SLIDE_WIDTH_PX * PX_TO_INCH)
PPTX_HEIGHT = Inches(SLIDE_HEIGHT_PX * PX_TO_INCH)

# Font mapping for better typography preservation
FONT_MAPPING = {
    'helvetica': 'Helvetica',
    'arial': 'Arial',
    'times new roman': 'Times New Roman',
    'georgia': 'Georgia',
    'verdana': 'Verdana',
    'tahoma': 'Tahoma',
    'trebuchet ms': 'Trebuchet MS',
    'garamond': 'Garamond',
    'comic sans ms': 'Comic Sans MS',
    'impact': 'Impact',
    'lucida sans unicode': 'Lucida Sans Unicode',
    'palatino linotype': 'Palatino Linotype',
    'calibri': 'Calibri',
    'cambria': 'Cambria',
    'candara': 'Candara',
    'consolas': 'Consolas',
    'corbel': 'Corbel',
    'franklin gothic medium': 'Franklin Gothic Medium',
    'gabriola': 'Gabriola',
    'geneva': 'Geneva',
    'lucida grande': 'Lucida Grande',
    'segoe ui': 'Segoe UI',
    'trebuchet ms': 'Trebuchet MS'
}

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def css_rgb_to_pptx_color(css_color):
    """Converts 'rgb(r, g, b)' or 'rgba(r, g, b, a)' to pptx RGBColor. Returns None for transparent."""
    if not css_color or 'rgba(0, 0, 0, 0)' in css_color or 'transparent' in css_color:
        return None
    
    nums = re.findall(r'\d+', css_color)
    if len(nums) >= 3:
        return RGBColor(int(nums[0]), int(nums[1]), int(nums[2]))
    return RGBColor(0, 0, 0) # Fallback

def map_alignment(align_str):
    if 'center' in align_str: return PP_ALIGN.CENTER
    if 'right' in align_str: return PP_ALIGN.RIGHT
    if 'justify' in align_str: return PP_ALIGN.JUSTIFY
    return PP_ALIGN.LEFT

def optimize_image(image_path, max_width=1920, max_height=1080, quality=85):
    """Optimize image to reduce file size while maintaining quality"""
    try:
        with Image.open(image_path) as img:
            # Calculate new dimensions maintaining aspect ratio
            original_width, original_height = img.size
            if original_width > max_width or original_height > max_height:
                ratio = min(max_width/original_width, max_height/original_height)
                new_width = int(original_width * ratio)
                new_height = int(original_height * ratio)
                img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # Save with optimized quality
            img.save(image_path, optimize=True, quality=quality)
    except Exception as e:
        print(f"Warning: Could not optimize image {image_path}: {e}")

async def generate_screenshot_pptx(html_uri, output_path):
    print(f"Generating Screenshot PPTX: {output_path}")
    prs = Presentation()
    prs.slide_width = PPTX_WIDTH
    prs.slide_height = PPTX_HEIGHT
    blank_slide_layout = prs.slide_layouts[6]

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page(viewport={'width': SLIDE_WIDTH_PX, 'height': SLIDE_HEIGHT_PX})
        
        print(f"Loading {html_uri}...")
        await page.goto(html_uri)
        
        try:
            await page.wait_for_load_state('networkidle', timeout=5000)
        except:
            pass

        slides = await page.query_selector_all('.slide')
        print(f"Found {len(slides)} slides.")

        for i, slide in enumerate(slides):
            screenshot_path = f"temp_slide_{i}.png"
            await slide.screenshot(path=screenshot_path)
            
            # Optimize the image before adding to presentation
            optimize_image(screenshot_path)
            
            pptx_slide = prs.slides.add_slide(blank_slide_layout)
            pptx_slide.shapes.add_picture(
                screenshot_path, 
                0, 0, 
                width=prs.slide_width, 
                height=prs.slide_height
            )
            if os.path.exists(screenshot_path):
                os.remove(screenshot_path)

        await browser.close()

    prs.save(output_path)
    print(f"Saved {output_path}")


async def generate_editable_pptx(html_uri, output_path):
    print(f"Generating Editable PPTX: {output_path}")
    prs = Presentation()
    prs.slide_width = PPTX_WIDTH
    prs.slide_height = PPTX_HEIGHT
    blank_layout = prs.slide_layouts[6]

    async with async_playwright() as p:
        browser = await p.chromium.launch(headless=True)
        page = await browser.new_page(viewport={'width': SLIDE_WIDTH_PX, 'height': SLIDE_HEIGHT_PX})
        
        await page.goto(html_uri)
        try:
            await page.wait_for_load_state('networkidle', timeout=5000)
        except:
            pass

        slide_handles = await page.query_selector_all('.slide')
        
        for i, slide_handle in enumerate(slide_handles):
            slide = prs.slides.add_slide(blank_layout)
            slide_box = await slide_handle.bounding_box()
            if not slide_box: continue
            slide_x, slide_y = slide_box['x'], slide_box['y']

            # --- LAYER 1: BACKGROUND SHAPES ---
            shapes_data = await slide_handle.evaluate("""(slide) => {
                const results = [];
                const slideRect = slide.getBoundingClientRect();
                const allEls = slide.querySelectorAll('*');
                
                allEls.forEach(el => {
                    const style = window.getComputedStyle(el);
                    const rect = el.getBoundingClientRect();
                    if (rect.width < 1 || rect.height < 1 || style.display === 'none' || style.visibility === 'hidden') return;
                    
                    const hasBg = style.backgroundColor !== 'rgba(0, 0, 0, 0)' && style.backgroundColor !== 'transparent';
                    const hasBorder = style.borderWidth !== '0px' && style.borderStyle !== 'none' && style.borderColor !== 'rgba(0, 0, 0, 0)';
                    const isAccent = el.classList.contains('accent-bar');
                    const isSection = el.classList.contains('slide-section');

                    if (hasBg || hasBorder || isAccent || isSection) {
                        if (el.classList.contains('slide')) return;
                        results.push({
                            x: rect.x - slideRect.x,
                            y: rect.y - slideRect.y,
                            w: rect.width,
                            h: rect.height,
                            bg: style.backgroundColor,
                            border: style.borderColor,
                            borderWidth: parseFloat(style.borderLeftWidth) || 0,
                            isAccent: isAccent
                        });
                    }
                });
                return results;
            }""")

            for shape in shapes_data:
                x = Inches(shape['x'] * PX_TO_INCH)
                y = Inches(shape['y'] * PX_TO_INCH)
                w = Inches(shape['w'] * PX_TO_INCH)
                h = Inches(shape['h'] * PX_TO_INCH)
                
                sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
                
                fill_color = css_rgb_to_pptx_color(shape['bg'])
                if shape['isAccent']:
                    sp.fill.solid()
                    sp.fill.fore_color.rgb = RGBColor(243, 244, 246)
                    sp.line.fill.background()
                elif fill_color:
                    sp.fill.solid()
                    sp.fill.fore_color.rgb = fill_color
                else:
                    sp.fill.background()

                if shape['borderWidth'] > 0:
                    sp.line.color.rgb = css_rgb_to_pptx_color(shape['border'])
                    sp.line.width = Pt(shape['borderWidth'])
                else:
                    sp.line.fill.background()

            # --- LAYER 2: IMAGES ---
            # Extended selector list to include more HTML elements
            image_selectors = [
                '.viz-box', '.dashboard-placeholder', '.bi', '.corner-icon', '.fa', 
                'img', 'svg', 'picture', 'canvas', 'video', 'figure', 'figcaption',
                '.image', '.photo', '.graphic', '.icon', '.logo', '.avatar',
                '.thumbnail', '.poster', '.banner', '.header-image', '.footer-image'
            ]
            
            for selector in image_selectors:
                elements = await slide_handle.query_selector_all(selector)
                for el in elements:
                    if not await el.is_visible(): continue
                    
                    box = await el.bounding_box()
                    if not box or box['width'] < 1: continue
                    
                    rel_x = box['x'] - slide_x
                    rel_y = box['y'] - slide_y
                    
                    screenshot_path = f"temp_img_{i}_{int(rel_x)}_{int(rel_y)}.png"
                    try:
                        await el.screenshot(path=screenshot_path)
                        
                        # Optimize the image before adding to presentation
                        optimize_image(screenshot_path)
                        
                        slide.shapes.add_picture(
                            screenshot_path, 
                            Inches(rel_x * PX_TO_INCH), 
                            Inches(rel_y * PX_TO_INCH),
                            width=Inches(box['width'] * PX_TO_INCH),
                            height=Inches(box['height'] * PX_TO_INCH)
                        )
                        if os.path.exists(screenshot_path):
                            os.remove(screenshot_path)
                    except:
                        pass

            # --- LAYER 3: TEXT ---
            # Enhanced text extraction with support for more HTML elements and nested structures
            text_data = await slide_handle.evaluate("""(slide) => {
                const results = [];
                const slideRect = slide.getBoundingClientRect();
                
                // Function to recursively extract text content while preserving structure
                function extractTextContent(element) {
                    let textContent = '';
                    const directTextNodes = Array.from(element.childNodes).filter(node => 
                        node.nodeType === 3 && node.nodeValue.trim().length > 0
                    );
                    
                    directTextNodes.forEach(node => {
                        textContent += node.nodeValue;
                    });
                    
                    return textContent.trim();
                }
                
                // Function to check if element has direct text content
                function hasDirectText(el) {
                    return Array.from(el.childNodes).some(node => 
                        node.nodeType === 3 && node.nodeValue.trim().length > 0
                    );
                }
                
                // Function to get all text content including nested elements
                function getAllTextContent(element) {
                    let text = '';
                    const walker = document.createTreeWalker(
                        element,
                        NodeFilter.SHOW_TEXT,
                        null,
                        false
                    );
                    
                    let node;
                    while (node = walker.nextNode()) {
                        text += node.nodeValue;
                    }
                    return text.trim();
                }
                
                // More comprehensive selector for various text elements
                const textSelectors = [
                    'h1', 'h2', 'h3', 'h4', 'h5', 'h6',
                    'p', 'span', 'div', 'a', 'strong', 'b', 'em', 'i', 'u', 's', 'sub', 'sup',
                    'li', 'td', 'th', 'caption', 'legend', 'label', 'button', 'input[type="button"]',
                    '.text', '.content', '.title', '.subtitle', '.headline', '.subheadline',
                    '.paragraph', '.description', '.note', '.caption', '.quote', '.blockquote',
                    '.highlight', '.emphasis', '.important', '.warning', '.alert', '.info',
                    '.header', '.footer', '.sidebar', '.nav', '.menu', '.breadcrumb',
                    '.tag', '.badge', '.chip', '.tooltip', '.popover', '.modal',
                    '.card', '.panel', '.box', '.section', '.container', '.wrapper'
                ];
                
                textSelectors.forEach(selector => {
                    const elements = slide.querySelectorAll(selector);
                    elements.forEach(el => {
                        if (el.closest('.viz-box') || el.closest('.dashboard-placeholder')) return;
                        if (el.tagName === 'I' || el.classList.contains('bi')) return;

                        if (hasDirectText(el)) {
                            const style = window.getComputedStyle(el);
                            const rect = el.getBoundingClientRect();
                            
                            // Get text content preserving structure
                            const textContent = extractTextContent(el);
                            
                            // Extract href for hyperlinks
                            let href = null;
                            if (el.tagName === 'A' && el.href) {
                                href = el.href;
                            }
                            
                            // Check if this element is inside an anchor tag
                            const parentAnchor = el.closest('a');
                            if (parentAnchor && !href) {
                                href = parentAnchor.href;
                            }
                            
                            results.push({
                                text: textContent, 
                                tagName: el.tagName,
                                x: rect.x - slideRect.x,
                                y: rect.y - slideRect.y,
                                w: rect.width,
                                h: rect.height,
                                color: style.color,
                                fontSize: style.fontSize,
                                fontFamily: style.fontFamily,
                                fontWeight: style.fontWeight,
                                textAlign: style.textAlign,
                                textTransform: style.textTransform,
                                textDecoration: style.textDecoration,
                                href: href, // Add hyperlink support
                                isNested: el.parentElement !== slide // Flag if nested
                            });
                        }
                    });
                });
                
                // Also check for any other elements that might contain text
                const allEls = slide.querySelectorAll('*');
                allEls.forEach(el => {
                    if (el.closest('.viz-box') || el.closest('.dashboard-placeholder')) return;
                    if (el.tagName === 'I' || el.classList.contains('bi')) return;
                    if (el.tagName === 'IMG' || el.tagName === 'SVG' || el.tagName === 'CANVAS') return;
                    
                    // Check if this element is already processed by the specific selectors
                    const isAlreadyProcessed = results.some(r => {
                        const existingEl = document.elementFromPoint(
                            r.x + slideRect.x + r.w/2, 
                            r.y + slideRect.y + r.h/2
                        );
                        return existingEl === el || el.contains(existingEl);
                    });
                    
                    if (!isAlreadyProcessed && hasDirectText(el)) {
                        const style = window.getComputedStyle(el);
                        const rect = el.getBoundingClientRect();
                        const textContent = extractTextContent(el);
                        
                        // Extract href for hyperlinks
                        let href = null;
                        if (el.tagName === 'A' && el.href) {
                            href = el.href;
                        }
                        
                        // Check if this element is inside an anchor tag
                        const parentAnchor = el.closest('a');
                        if (parentAnchor && !href) {
                            href = parentAnchor.href;
                        }
                        
                        results.push({
                            text: textContent, 
                            tagName: el.tagName,
                            x: rect.x - slideRect.x,
                            y: rect.y - slideRect.y,
                            w: rect.width,
                            h: rect.height,
                            color: style.color,
                            fontSize: style.fontSize,
                            fontFamily: style.fontFamily,
                            fontWeight: style.fontWeight,
                            textAlign: style.textAlign,
                            textTransform: style.textTransform,
                            textDecoration: style.textDecoration,
                            href: href, // Add hyperlink support
                            isNested: el.parentElement !== slide
                        });
                    }
                });
                
                return results;
            }""")

            for txt in text_data:
                content = txt['text'].strip()
                if not content: continue

                tx = Inches(txt['x'] * PX_TO_INCH)
                ty = Inches(txt['y'] * PX_TO_INCH)
                tw = Inches(txt['w'] * PX_TO_INCH)
                th = Inches(txt['h'] * PX_TO_INCH)

                # Better text wrapping and overflow handling
                # Ensure minimum dimensions for text boxes
                if tw < Inches(0.5):
                    tw = Inches(0.5)
                if th < Inches(0.25):
                    th = Inches(0.25)

                textbox = slide.shapes.add_textbox(tx, ty, tw, th)
                tf = textbox.text_frame
                tf.word_wrap = True
                tf.auto_size = False  # Disable auto-size for better control
                
                # Better text overflow handling
                tf.fit_text = True  # Fit text to shape
                tf.margin_top = Pt(2)
                tf.margin_bottom = Pt(2)
                tf.margin_left = Pt(2)
                tf.margin_right = Pt(2)
                
                p_node = tf.paragraphs[0]
                
                # Handle special list items
                if txt['tagName'] == 'LI' and not content.startswith("■"):
                    content = "■ " + content

                if txt['textTransform'] == 'uppercase':
                    content = content.upper()

                p_node.text = content
                p_node.alignment = map_alignment(txt['textAlign'])
                
                run = p_node.runs[0]
                rgb = css_rgb_to_pptx_color(txt['color'])
                if rgb: 
                    run.font.color.rgb = rgb
                
                # Handle font size with better precision
                size_match = re.match(r'([\d.]+)px', txt['fontSize'])
                if size_match:
                    # Convert px to points (1px = 0.75pt)
                    px_size = float(size_match.group(1))
                    # Apply scaling factor for better visual match
                    scaled_size = px_size * 0.75
                    run.font.size = Pt(scaled_size) 
                
                # Font mapping for better typography preservation
                font_family = txt['fontFamily'].lower()
                matched_font = None
                for key, value in FONT_MAPPING.items():
                    if key in font_family:
                        matched_font = value
                        break
                if matched_font:
                    run.font.name = matched_font
                else:
                    # Fallback to common fonts
                    if 'helvetica' in font_family: run.font.name = 'Helvetica'
                    elif 'arial' in font_family: run.font.name = 'Arial'
                    elif 'times' in font_family: run.font.name = 'Times New Roman'
                    else: run.font.name = 'Calibri'  # Default to Calibri
                
                # Handle font weight
                if 'bold' in str(txt['fontWeight']) or (str(txt['fontWeight']).isdigit() and int(txt['fontWeight']) >= 600):
                    run.font.bold = True
                
                # Handle text decoration (underline, strikethrough)
                if txt['textDecoration'] and 'underline' in txt['textDecoration']:
                    run.font.underline = True
                if txt['textDecoration'] and ('line-through' in txt['textDecoration'] or 'strikethrough' in txt['textDecoration']):
                    run.font.strike = True
                
                # Add hyperlink support
                if txt['href']:
                    try:
                        run.hyperlink.address = txt['href']
                    except:
                        pass  # If hyperlink fails, continue without it

        await browser.close()

    prs.save(output_path)
    print(f"Saved {output_path}")


async def main():
    parser = argparse.ArgumentParser(description="Convert HTML presentation to PPTX (Screenshot & Editable)")
    parser.add_argument("input_file", help="Path to the HTML file")
    args = parser.parse_args()

    input_path = os.path.abspath(args.input_file)
    if not os.path.exists(input_path):
        print(f"Error: File {input_path} not found.")
        sys.exit(1)

    base_name = os.path.splitext(input_path)[0]
    # Use the current working directory for output, not the source directory
    cwd = os.getcwd()
    file_name = os.path.basename(base_name)
    
    output_screenshot = os.path.join(cwd, f"{file_name}.pptx")
    output_editable = os.path.join(cwd, f"{file_name}_editable.pptx")
    
    html_uri = f"file://{input_path}"

    await generate_screenshot_pptx(html_uri, output_screenshot)
    await generate_editable_pptx(html_uri, output_editable)
    
    print("Done.")


if __name__ == "__main__":
    asyncio.run(main())